"""文档解析:PDF / PPT / Markdown / 文本 / 图片。"""
from __future__ import annotations

from pathlib import Path

# 支持的离线文档类型
SUPPORTED_TYPES = {
    "pdf", "pptx", "ppt", "md", "markdown", "txt",
    "png", "jpg", "jpeg", "gif", "webp",
}

# 图片类型(当前不做 OCR,仅存原图,不向量化)
IMAGE_TYPES = {"png", "jpg", "jpeg", "gif", "webp"}


def parse_pdf(path: str) -> str:
    """解析 PDF,逐页提取文本。"""
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = [(page.extract_text() or "") for page in reader.pages]
    return "\n".join(pages)


def parse_pptx(path: str) -> str:
    """解析 PPT/PPTX,提取所有幻灯片文本。"""
    from pptx import Presentation

    prs = Presentation(path)
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        lines.append(line)
    return "\n".join(lines)


def parse_text(path: str) -> str:
    """解析纯文本 / Markdown。"""
    return Path(path).read_text(encoding="utf-8", errors="ignore")


def parse_file(path: str, file_type: str) -> str:
    """统一解析入口,返回纯文本;图片返回空字符串(不解析)。"""
    ft = (file_type or "").lower().lstrip(".")
    if ft == "pdf":
        return parse_pdf(path)
    if ft in ("pptx", "ppt"):
        return parse_pptx(path)
    if ft in ("md", "markdown", "txt"):
        return parse_text(path)
    return ""  # 图片等暂不解析


def is_image(file_type: str) -> bool:
    return (file_type or "").lower().lstrip(".") in IMAGE_TYPES
