"""生成演示测试数据:产品部、销售部 各 3 份企业知识文档(MD + PDF + PPT)。

输出目录:
    data/demo_data/产品/  产品部-产品需求与开发流程.md
                         产品部-星云智能音箱X2产品规格书.pdf
                         产品部-新品X2产品培训.pptx
    data/demo_data/销售/  销售部-销售管理制度.md
                         销售部-产品定价与折扣政策.pdf
                         销售部-新品X2上市销售培训.pptx

用法:
    uv run python scripts/generate_demo_data.py
"""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

OUT = Path("data/demo_data")
HEI = "C:/Windows/Fonts/simhei.ttf"  # 黑体(中文字体)

# ============================= 产品部 =============================

PRODUCT_MD = """# 星云科技 产品部 — 产品需求与开发流程

## 一、需求提交与评审

- 需求来源:客户反馈、销售提报、内部产品规划、竞品分析。
- 所有需求须在 OA 系统填写《产品需求单》,包含:背景、目标用户、功能描述、验收标准。
- 每周三上午召开需求评审会,由产品经理、研发负责人、测试负责人共同评审。
- 需求优先级分为四级:
  - P0:紧急缺陷/合规问题,立即处理;
  - P1:高优先级,当个迭代必须完成;
  - P2:中优先级,排入后续迭代;
  - P3:低优先级,放入需求池待评估。

## 二、产品开发流程

开发分为八个阶段,每个阶段有明确产出物:

1. 需求评审 → 《需求评审纪要》
2. 原型设计 → 《交互原型稿》
3. UI 评审 → 《UI 视觉稿》
4. 技术评审 → 《技术方案文档》
5. 开发实现 → 代码 + 单元测试
6. 联调测试 → 冒烟测试、功能测试、回归测试
7. 灰度发布 → 内部员工 5% 灰度,观察 3 天
8. 全量上线 → 发布公告 + 回滚预案

## 三、版本命名规则

- 版本号格式:主版本.次版本.修订号,例如 1.0.0。
- 主版本:重大功能或架构变更时递增;
- 次版本:新增功能时递增;
- 修订号:缺陷修复时递增。
- 版本发布前须通过全部测试用例,并在《发布检查单》上逐项确认。

## 四、需求变更流程

- 需求变更须填写《需求变更单》,说明变更原因与影响范围;
- 变更影响评估由产品经理发起,研发与测试确认工作量;
- 涉及上线时间的变更须产品总监审批;
- 已进入开发阶段的需求变更,须评估返工成本后再决策。

## 五、产品文档规范

- PRD 结构:背景、目标、用户故事、功能详述、验收标准、数据埋点。
- 验收标准须可量化,例如"首屏加载时间小于 2 秒"。
- 文档统一存放于知识库"产品文档"目录,按版本归档。
"""

PRODUCT_PDF_TITLE = "星云科技 产品部 — 星云智能音箱 X2 产品规格书"
PRODUCT_PDF_SECTIONS = [
    ("一、产品概述", [
        "星云智能音箱 X2 是星云科技 2026 年旗舰智能音箱,搭载第二代自研语音大模型。",
        "定位:高端智能家居控制中枢,支持全屋智能设备联动。",
    ]),
    ("二、核心参数", [
        "型号:X2;官方售价:699 元人民币。",
        "尺寸:130mm × 130mm × 200mm;重量:950 克。",
        "扬声器:2.5 英寸全频单元 + 被动低音辐射器。",
        "麦克风:6 阵列麦克风,支持 8 米远场拾音。",
        "无线连接:Wi-Fi 6E、蓝牙 5.4;支持 2.4GHz / 5GHz / 6GHz 频段。",
        "颜色:星空灰、晨曦金;电源:DC 12V 2A 适配器。",
    ]),
    ("三、软件功能", [
        "语音助手:自然语言对话、天气查询、闹钟提醒、日程管理。",
        "音乐播放:接入主流音乐平台,支持无损音质与多房间播放。",
        "智能家居:支持米家、HomeKit 生态,可语音控制灯光、窗帘、空调、门锁。",
        "多设备组网:支持四台音箱组成立体声/多房间系统。",
        "固件升级:支持 OTA 自动升级,升级包约 120MB。",
    ]),
    ("四、认证与合规", [
        "已通过 3C 认证、SRRC 无线电发射设备型号核准。",
        "符合 RoHS 环保指令;包装采用可回收材料。",
    ]),
    ("五、包装清单与保修", [
        "包装清单:主机 ×1、电源适配器 ×1、快速指南 ×1、保修卡 ×1。",
        "保修政策:整机保修 1 年,主要部件(扬声器/麦克风)保修 3 年。",
        "退换货:7 天无理由退货,15 天内质量问题换新。",
    ]),
    ("六、常见问题", [
        "无法连接 Wi-Fi:确认路由器支持 2.4GHz 频段,并靠近音箱后重试。",
        "语音唤醒不灵敏:在 App 中重新校准唤醒词,并避免遮挡麦克风。",
        "固件升级失败:保持设备通电并连接网络,在 App 中重新触发升级。",
    ]),
]

# ============================= 销售部 =============================

SALES_MD = """# 星云科技 销售部 — 销售管理制度

## 一、客户分级管理

- A 级客户:年采购额 100 万元以上,由销售总监直接跟进;
- B 级客户:年采购额 30 万至 100 万元,由资深销售跟进;
- C 级客户:年采购额 30 万元以下,由普通销售跟进。
- 客户信息须在签单后 48 小时内录入 CRM 系统。

## 二、报价与折扣权限

- 报价单有效期 30 天,逾期须重新审批。
- 折扣权限分级:
  - 销售员:最高 5% 折扣;
  - 销售主管:最高 10% 折扣;
  - 销售总监:最高 15% 折扣;
  - 超过 15% 的折扣须总经理特批。
- 所有报价须在 CRM 中留痕,禁止口头承诺价格。

## 三、合同管理

- 合同金额 10 万元以上须经法务审核后签订;
- 合同盖章须走 OA 用印流程,禁止私自盖章;
- 合同签订后 3 个工作日内归档至合同管理系统。

## 四、提成与绩效

- 提成按实际回款计提:基础提成比例 3%,超额完成季度目标部分提成 5%;
- 回款周期超过 90 天的订单,提成比例减半;
- 客户退货订单,对应提成从下月发放中扣除。

## 五、客户拜访要求

- A 级客户每月至少拜访 1 次,B 级客户每季度至少 1 次;
- 拜访记录须在 48 小时内录入 CRM;
- 出差拜访须提前在 OA 提交出差申请并经主管审批。
"""

SALES_PDF_TITLE = "星云科技 销售部 — 产品定价与折扣政策"
SALES_PDF_SECTIONS = [
    ("一、产品价目表", [
        "星云智能音箱 X1:官方零售价 499 元。",
        "星云智能音箱 X2:官方零售价 699 元。",
        "智能网关 Z1:官方零售价 299 元。",
        "智能传感器套装(3 件装):官方零售价 199 元。",
    ]),
    ("二、渠道折扣政策", [
        "一级经销商:进货价按零售价 7 折(30% 折扣)。",
        "二级经销商:进货价按零售价 8 折(20% 折扣)。",
        "单笔采购 50 台及以上,额外享受 5% 折上折。",
    ]),
    ("三、大客户与政企折扣", [
        "年采购额 100 万元以上大客户,可申请 15% 以内阶梯折扣。",
        "政企项目采购须单独报价,由销售总监审批。",
    ]),
    ("四、促销政策", [
        "官方大促(双 11、618):全渠道直降 10%。",
        "新品上市首月:购买 X2 赠送智能传感器套装。",
        "以旧换新:旧款 X1 抵扣 100 元。",
    ]),
    ("五、退换与价格保护", [
        "7 天无理由退货,15 天内质量问题换新。",
        "调价前 30 天书面通知经销商,并提供价格保护。",
        "退货须保持包装完整、配件齐全,运费由责任方承担。",
    ]),
]

# ============================= PPT 内容 =============================

PRODUCT_PPT = [
    ("星云智能音箱 X2 — 产品培训", "产品部 | 2026 年 8 月"),
    ("产品定位", [
        "旗舰级智能音箱,定位高端智能家居控制中枢",
        "目标用户:品质家庭、智能家居深度用户",
        "差异化:8 米远场拾音 + Wi-Fi 6E + 全屋联动",
    ]),
    ("核心卖点", [
        "卖点一:8 米远场拾音,厨房也能轻松唤醒",
        "卖点二:Wi-Fi 6E 三频段,网络更稳定",
        "卖点三:支持四台组网,打造全屋多房间音乐",
        "卖点四:兼容米家与 HomeKit,一屏管全家",
    ]),
    ("核心参数", [
        "官方售价:699 元;尺寸 130×130×200mm;重量 950 克",
        "扬声器:2.5 英寸全频单元 + 被动低音辐射器",
        "麦克风:6 阵列,8 米远场;无线:Wi-Fi 6E / 蓝牙 5.4",
        "颜色:星空灰、晨曦金",
    ]),
    ("目标用户画像", [
        "人群一:35-45 岁品质家庭,注重音质与智能联动",
        "人群二:智能家居爱好者,已有米家/HomeKit 设备",
        "人群三:礼品场景,替换传统音箱的升级需求",
    ]),
    ("竞品对比", [
        "vs 竞品 A:远场拾音距离更长(8 米 vs 5 米)",
        "vs 竞品 B:支持四台组网,竞品仅支持双台",
        "vs 竞品 C:价格更低(699 vs 899),认证更全",
    ]),
    ("演示脚本", [
        "1. 唤醒演示:8 米外喊「你好星云」",
        "2. 音乐演示:播放无损音源,展示低音表现",
        "3. 联动演示:语音关闭窗帘、调节空调",
        "4. 组网演示:两台 X2 组立体声",
    ]),
    ("Q&A 与资料", [
        "常见问题见《X2 产品规格书》",
        "培训资料存放于知识库「产品文档」目录",
        "联系方式:产品部 - 王工(内线 8123)",
    ]),
]

SALES_PPT = [
    ("星云智能音箱 X2 — 上市销售培训", "销售部 | 2026 年 8 月"),
    ("销售目标", [
        "上市首月目标:全国销量 2 万台",
        "重点渠道:线上旗舰店 + 一级经销商",
        "重点场景:智能家居套餐、以旧换新",
    ]),
    ("目标客户", [
        "A 类:智能家居渠道商,采购 50 台以上享折上折",
        "B 类:政企礼品采购,关注品牌与认证",
        "C 类:个人家庭用户,关注音质与价格",
    ]),
    ("产品卖点话术", [
        "痛点「音箱太远听不清」 → 卖点:8 米远场拾音",
        "痛点「家里设备各管各的」 → 卖点:米家/HomeKit 全屋联动",
        "痛点「网速卡顿」 → 卖点:Wi-Fi 6E 三频段",
    ]),
    ("价格与折扣", [
        "零售价 699 元;一级经销商 7 折、二级 8 折",
        "单笔 50 台以上额外 5% 折上折",
        "大促直降 10%,新品首月赠送传感器套装",
    ]),
    ("常见异议处理", [
        "异议「太贵了」 → 对比 X1 升级点与赠品价值",
        "异议「和竞品差不多」 → 突出 8 米拾音 + 四台组网",
        "异议「再等等」 → 强调首月赠品限量、以旧换新抵扣 100 元",
    ]),
    ("跟进与记录", [
        "客户信息 48 小时内录入 CRM",
        "报价单有效期 30 天,超期重新审批",
        "A 级客户每月拜访 1 次,拜访记录 48 小时内录入",
    ]),
    ("Q&A", [
        "产品参数疑问 → 查《X2 产品规格书》",
        "折扣审批 → 按《销售管理制度》权限逐级上报",
        "培训资料 → 知识库「销售文档」目录",
    ]),
]


# ============================= 生成实现 =============================


def write_md(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    print(f"  [MD ] {path.relative_to('data')}")


def write_pdf(path: Path, title: str, sections: list) -> None:
    from fpdf import FPDF

    path.parent.mkdir(parents=True, exist_ok=True)
    pdf = FPDF(format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_font("hei", "", HEI)
    pdf.add_page()

    # 标题
    pdf.set_font("hei", size=16)
    pdf.multi_cell(0, 10, title, align="C")
    pdf.ln(2)

    for sec_title, bullets in sections:
        # 章节标题
        pdf.set_font("hei", size=13)
        pdf.multi_cell(0, 9, sec_title)
        pdf.ln(1)
        # 条目
        pdf.set_font("hei", size=11)
        for b in bullets:
            pdf.multi_cell(0, 7, f"- {b}")
            pdf.ln(0.5)
        pdf.ln(2)

    pdf.output(str(path))
    print(f"  [PDF] {path.relative_to('data')}")


def write_pptx(path: Path, slides: list) -> None:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = 12192000  # 16:9
    prs.slide_height = 6858000
    for idx, (title, body) in enumerate(slides):
        layout = prs.slide_layouts[1] if idx > 0 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if isinstance(body, list):
            tf = slide.placeholders[1].text_frame
            tf.text = body[0]
            for line in body[1:]:
                p = tf.add_paragraph()
                p.text = line
        else:
            slide.placeholders[1].text = body
        # 设置中文字体(同时设置 Latin 与 East Asian)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(18 if idx == 0 else 16)
                    rPr = run._r.get_or_add_rPr()
                    ea = rPr.find(qn("a:ea"))
                    if ea is None:
                        ea = rPr.makeelement(qn("a:ea"), {})
                        rPr.append(ea)
                    ea.set("typeface", "微软雅黑")
    prs.save(str(path))
    print(f"  [PPT] {path.relative_to('data')}")


def main() -> None:
    print("生成演示数据:")
    product_dir = OUT / "产品"
    sales_dir = OUT / "销售"

    write_md(product_dir / "产品部-产品需求与开发流程.md", PRODUCT_MD)
    write_pdf(product_dir / "产品部-星云智能音箱X2产品规格书.pdf", PRODUCT_PDF_TITLE, PRODUCT_PDF_SECTIONS)
    write_pptx(product_dir / "产品部-新品X2产品培训.pptx", PRODUCT_PPT)

    write_md(sales_dir / "销售部-销售管理制度.md", SALES_MD)
    write_pdf(sales_dir / "销售部-产品定价与折扣政策.pdf", SALES_PDF_TITLE, SALES_PDF_SECTIONS)
    write_pptx(sales_dir / "销售部-新品X2上市销售培训.pptx", SALES_PPT)

    print("\n完成!文件位于 data/demo_data/ 下,可直接在前端上传到对应部门的知识库。")


if __name__ == "__main__":
    main()
